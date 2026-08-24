#!/usr/bin/env python3
"""
make_papers_deck.py - the two reference papers, figure by figure, against our results.

One slide per figure. The slide shows OUR corresponding figure; the notes describe
what their figure does, how ours compares, and whether we reproduced the method,
departed from it, or cannot do it at all. Their figures are not reproduced here -
described only, with the section or panel named so any claim can be checked.

    Norman-Haignere et al. (2019), bioRxiv 696161 - ECoG, 13 patients, 165 sounds,
    sparse component model. Main figures 1-4.

    Castellucci et al. (2026), Cell Reports 45:116783 - ECoG/sEEG, 22 patients,
    question-answer speech task. Main figures 1-4 plus S1-S5.

    python make_papers_deck.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent
OUT = ROOT / "outputs" / "review_deck"

INK = RGBColor(0x1B, 0x23, 0x2C)
MUTED = RGBColor(0x68, 0x72, 0x7D)
NH = RGBColor(0x7A, 0x4E, 0x9B)          # Norman-Haignere
CA = RGBColor(0x4A, 0x6F, 0xA5)          # Castellucci
GREEN = RGBColor(0x1B, 0x78, 0x37)
RED = RGBColor(0xC1, 0x12, 0x1F)

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)

C7 = "02_FBM_Clustering/outputs/clustering/comparison/C7_separation.png"
D1 = ("02_FBM_Clustering/outputs/clustering/cnmf/concat_hg/runs/20260818_112939/"
      "decomposition_drivers/D1_membership.png")
D2 = ("02_FBM_Clustering/outputs/clustering/cnmf/concat_hg/runs/20260818_112939/"
      "decomposition_drivers/D2_drivers.png")
D3 = "02_FBM_Clustering/outputs/clustering/decomposition/concat_hg/D3_graded_decomposition.png"
G1 = "02_FBM_Clustering/outputs/clustering/comparison/G1_silhouette_by_gate.png"
G2 = "02_FBM_Clustering/outputs/clustering/comparison/G2_drivers_by_gate.png"
G3 = "02_FBM_Clustering/outputs/clustering/gate_examples/G3_gate_borderline_examples.png"
A1 = "02_FBM_Clustering/outputs/clustering/component_anatomy/A1_component_anatomy.png"
B1 = "02_FBM_Clustering/outputs/clustering/overfactor_stability/boot_cnmf_concat_hg_K10.png"

SLIDES = [
    ("Two papers, figure by figure",
     "What each one does, and what we took from it",
     None, INK,
     """WHY THESE TWO
Norman-Haignere et al. 2019 — ECoG, 13 patients, 165 sounds, sparse component model.
The methodological template for graded components and for filtering them.
Castellucci et al. 2026 (Cell Reports 45:116783) — ECoG/sEEG, 22 patients, a
question-answer speech task. The closest published analysis to ours in modality, task
and cohort size.

HOW TO READ THESE SLIDES
Each slide shows OUR figure. The notes describe THEIR figure and the comparison.
Their figures are not reproduced — described only, with the panel named.

THE ONE-LINE SUMMARY
We reproduce Norman-Haignere's component-filtering discipline and Castellucci's
anatomical statistics, we depart from both on electrode selection, and we add two
controls neither applies."""),

    ("The two papers at a glance",
     "Where they agree, and where we sit",
     None, INK,
     """                         NORMAN-HAIGNERE 2019      CASTELLUCCI 2026        US
Cohort                   13 patients               22 patients             27
Electrodes kept          271                       1296 of 1784 (72.6%)    1266 of 2946 (43%)
Selection criterion      split-half r > 0.2        significance vs a       AMPLITUDE
                                                   shuffled null           threshold
Feature space            165 sounds x 3 s          4 epochs concatenated,  3 conditions
                                                   z-scored, PCA -> 3 PCs  x 300 bins, 900-D
Model                    sparse NMF + learned      plain NMF               convex NMF
                         smoothing kernel, L1
Choosing K               cell-wise CV, noise-      RMSE/VAF elbow          bi-cross-validation
                         corrected
Component stability      2 filters (init 1000x,    NONE                    bootstrap + CI +
                         subject weight)                                   null + init + patient
Hard partition?          never attempted           attempted, FAILED       attempted, FAILED
Anatomy                  descriptive only          permutation-tested,     tested, two nulls
                                                   uncorrected             + correction

THE PATTERN: both select on REPRODUCIBILITY. We select on amplitude. That is the one
place we are behind both, and it is what the split-half work fixes."""),

    ("NH Fig 1A–D — the setup and the data matrix",
     "Their design ↔ our concat_hg matrix",
     D3, NH,
     """THEIR FIGURE
1A — the sound set: 165 everyday sounds, 2 s each.
1B — an electrode map coloured by split-half reliability of broadband gamma (70–140 Hz).
1C — their earlier fMRI speech/music selectivity maps with sound-responsive electrodes
     overlaid, as a prior expectation about where things should be.
1D — the decomposition schematic: rows are electrodes, columns are the concatenated
     response timecourse across all 165 sounds; the matrix is approximated as a product
     of an electrode-weight matrix and a component-response matrix.

OURS
Same shape of object. Our matrix is 1266 electrodes x 900 features, where the 900 are
three conditions x 300 time-normalised bins (audio, picture, reading), GO cue at bin 150.

WHAT DIFFERS
Their columns are STIMULI (165 sounds); ours are TIME within three task conditions.
So their components are tuning profiles across a sound set, ours are temporal response
shapes. Both are "a response profile per component, a weight per electrode".

REPRODUCED: the matrix formulation and the graded weight interpretation."""),

    ("NH Fig 1B — split-half reliability as the electrode filter",
     "Their selection criterion ↔ the gate we are replacing",
     G3, NH,
     """THEIR FIGURE
1B plots the split-half Pearson correlation of broadband gamma for every electrode.
Their 271 analysed electrodes are exactly the ones above r > 0.2. Stated in Results,
"Electrode decomposition".

They then REUSE the same quantity to noise-correct their variance explained (Fig 1E),
so one measurement buys both the selection and the denominator.

OURS — THE SLIDE
Our gate counts bins over a threshold: ≥2% of the 0–400 Hz cube above +2.2 dB, OR ≥4%
below −3.0 dB. 1323 of 3002 contacts pass.

THE COMPARISON IS THE ARGUMENT
Their criterion asks: does this response REPEAT?
Ours asks: is this response LARGE?
Those come apart. EL033 aH_R5 has 774 bins over threshold and is kept; aH_R2 — same
shank, same patient — has 773 and is discarded. One bin in 38,700.

STATUS: adopting theirs. Stage 01 now writes odd/even half-cubes; validated (real
response r = +0.908, pure noise r = +0.037); needs 140 re-run."""),

    ("NH Fig 1E — how many components?",
     "Their cross-validation ↔ our bi-cross-validation",
     D3, NH,
     """THEIR FIGURE
1E is the component-count criterion. The data matrix is divided into CELLS — one cell is
one electrode's response to one sound. They train on a random 80% of cells and predict
the remaining 20%, then plot squared test correlation against the number of components.

Two details worth copying:
  - the correlation is NOISE-CORRECTED using the electrodes' test-retest reliability, so
    it measures explainable rather than total variance
  - error bars come from bootstrapping ACROSS SUBJECTS, not across electrodes

OURS
Bi-cross-validation — held out in rows AND columns. Rises from 0.386 at K=2 to a peak
near K=8, with K=7–10 inside one standard error.

WHAT WE HAVE THAT THEY DO NOT: holding out columns as well as rows.
WHAT THEY HAVE THAT WE DO NOT: the noise correction, because it needs split-half
reliability — which the half-cubes now unlock. This is the second thing that work buys."""),

    ("NH Fig 1F & 1G — filtering the components",
     "Their two filters ↔ our five (notebook 238)",
     B1, NH,
     """THEIR FIGURES
1F — average weight of each component in each subject, normalised to sum to 1 across
subjects. A component whose maximum exceeds 0.5 is judged subject-specific. Components
14, 16, 18, 19 and 20 died this way.
1G — the decomposition was run 1000 times from different random initialisations; they
plot the median correlation between the best solution and the next 99 best, and keep
components above 0.9.

THE CRITICAL DETAIL: they fitted 20 components and interpreted 14. The MODEL STAYED AT
20. They never refit at 14, and never claimed 14 explain the data. The filter restricts
interpretation, not the model.

OURS — notebook 238 keeps that discipline and adds three filters:
  bootstrap Jaccard (the one that bites here), a confidence interval on it, and a
  calibrated null.
Their subject filter never bites on our data — 27 patients makes a >0.5 single-patient
share nearly unreachable; our maximum is 0.389.

REPRODUCED: the over-factor-then-filter strategy, and the model-stays-at-K discipline."""),

    ("NH Fig 2 — component responses and where they sit",
     "Their descriptive anatomy ↔ our tested anatomy (FIG C.11)",
     A1, NH,
     """THEIR FIGURE
2A — each component's response to all 165 sounds as a raster, sounds grouped into 12
categories, with category averages below and the time-averaged response at the right.
2B — anatomical maps of the electrode weights for each component, projected to the
cortical surface and aligned to fsaverage.

THE SENTENCE THAT MATTERS: they state plainly that electrode anatomy played no role in
the component analysis, and Fig 2B is presented as a DESCRIPTION of where the weights
fall — not a statistical test. There is no permutation test on anatomy anywhere in the
paper.

OURS — FIG C.11 turns that into a test. For every 1 cm voxel holding ≥3 of our 1029
localised contacts, is the median loading higher than chance allows? 14,927 voxels,
10,000 permutations.

Result: c1 (708 voxels FDR, 51% somatomotor) and c5 (242, 81% default) survive; c3 and
c6 localise nowhere.

THIS IS WHERE WE GO BEYOND THEM — they show maps, we test them."""),

    ("NH Fig 3 — model-matched sound synthesis",
     "No analogue in our data, and worth saying why",
     None, MUTED,
     """THEIR FIGURE
3 tests whether the selective components can be explained by generic acoustic
representations. They synthesise sounds matched to natural ones in modulation statistics
but otherwise unstructured, and compare the components' responses to natural versus
modulation-matched pairs. The selective components respond much less to the matched
sounds — so the selectivity is not explained by the acoustic model.

WHY WE HAVE NOTHING LIKE IT
This is a STIMULUS-CONTROL experiment. It requires generating new stimuli and
re-recording. Our design has three fixed conditions and no synthesised control set.

WHAT IT WOULD CORRESPOND TO
The equivalent question for us is whether a component's response is explained by
low-level properties of the stimulus rather than by language processing — and the
nearest thing we have is the LanA correlation map (FIG 9), which shows language-network
proximity predicts MORE high gamma while hearing speech and LESS during visual
encoding. That is a dissociation, but it is correlational and not a control condition.

SAY THIS IF ASKED: we cannot rule out acoustic confounds by their method. We can only
show the anatomy dissociates."""),

    ("NH Fig 4 — individual electrodes, no model",
     "Their per-electrode check ↔ our FIG C.5",
     D1, NH,
     """THEIR FIGURE
4 steps outside the component model entirely. They select individual speech-, music- and
song-selective electrodes and show their raw responses, to demonstrate the selectivity
exists in single electrodes and is not manufactured by the decomposition.

OURS — FIG C.5 asks the mirror-image question: does the model assign individual
electrodes cleanly? Median top weight 0.43 against 0.14 for an even split; only 34% have
a majority component; 236 electrodes (19%) have their top two within 0.05.

THE CONTRAST IS INSTRUCTIVE
They use single electrodes to CONFIRM the components are real.
We use single electrodes to show the LABELS are unreliable while the components are not
(FIG C.6: loading correlation and Cohen's d agree at r = 0.97–0.99).

Both land in the same place: the component is the trustworthy object, the per-electrode
assignment is not."""),

    ("Castellucci Fig 1 — four clusters, and a null that passed",
     "Their clustering ↔ our FIG C.7",
     C7, CA,
     """THEIR FIGURE
1A electrode coverage. 1B the task. 1C an example electrode with significant activation
and suppression marked. 1D variance explained per principal component — only the first
three exceed 5%. 1E silhouette and Calinski-Harabasz for four clustering methods over
K=2–10. 1F the distribution of PC coefficients per cluster.

They cluster on the coefficients of the FIRST 3 PCs, compare k-means, k-medoids,
Gaussian mixture and Ward, and pick k-means at K=4.

Crucially they ran a null: Parallel Analysis (Horn 1965) reruns the whole PCA-and-
clustering pipeline on shuffled responses, and the clustering disappears (Figs S2C–S2F).

OURS — FIG C.7 runs the same kind of test and our argmax partition does NOT clear it:
convex NMF scores z = +1.0 in dB. Scored in its own space it reaches z = +15.

THE LIKELY REASON, AND IT IS TESTABLE
They cluster in 3 dimensions. We cluster in 900. Silhouette is not space-free — that is
the same dependence that made our first version of C.7 wrong. Their K=4 and our K=7 are
not comparable numbers.

TO DO: reduce to a few PCs, re-run the separation test."""),

    ("Castellucci Fig 2 — the four response classes",
     "Their class #4 ↔ what our gate throws away",
     G1, CA,
     """THEIR FIGURE
2A group responses at four task alignment points. 2B individual electrode heterogeneity.
2C the class-level pattern. 2D the four classes with sizes. 2E when each class is active
relative to task structure.

THE FOUR CLASSES
  #1 sensory   — active after question onset, suppressed after offset   n=236 (18.2%)
  #2 planning  — active after critical information, before the answer   n=473 (36.5%)
  #3 motor     — active during and just before the answer               n=315 (24.3%)
  #4 SUPPRESSED — weakly modulated, generally suppressed                n=272 (21.0%)

CLASS #4 IS THE ONE TO POINT AT
A fifth of their modulated electrodes are weak-and-suppressed, and it survives as a
genuine, anatomically localised class.

OURS — that is the population our amplitude gate removes. FIG C.9 shows k-means and Ward
building a cluster 74% and 71% made of exactly those contacts once the gate is lifted,
and we read it as a failure mode.

THE HONEST CAVEAT: their class-4 electrodes passed a test for significant SUPPRESSION.
Ours passed no test — they were simply under a threshold. So the populations are not the
same, and ours could still be junk. The difference between them is precisely the
criterion we are moving to."""),

    ("Castellucci Fig 3 — classes localise to distinct substrates",
     "Their spatial statistics ↔ ours, and the control they run",
     A1, CA,
     """THEIR FIGURE
3A example participants coloured by class. 3B the proportion of neighbouring electrodes
(within 1 cm) sharing a class. 3C the same within participant, 3D pooled across
participants, each against 1000 label shuffles. 3E–3H the four classes on the cortical
surface, with voxels holding significantly more electrodes of that class than chance.

THEIR VOXEL TEST: MNI152 divided into 1 cm cubes with centres 1 mm apart; shuffle class
labels 1000x; a voxel is significant above the 95th percentile; voxels with ≤2
electrodes are not tested.

THE CONTROL WORTH COPYING: 3C versus 3D — they run the neighbour test BOTH within
participant and pooled, so "classes are anatomically organised" is separated from
"one patient's electrodes are near each other and share a class".

OURS — FIG C.11 uses their voxel machinery but adds that same within-patient control to
the WEIGHT test, which they only apply to the class test. We also correct for
multiplicity: 746 of our 14,927 voxels pass on noise alone.

CAVEAT THAT CUTS AGAINST US: their ECoG grids are dense and contiguous; our sEEG sits at
a median of 4 contacts per voxel, so their medians are far better estimated."""),

    ("Castellucci Fig 4 + S5 — NMF, and the wall we both hit",
     "Their pivot ↔ our FIG C.11",
     D2, CA,
     """THEIR FIGURE
4A a high-density grid participant. 4C activity flowing across regions in 300 ms windows.
4D and 4E the NMF component profiles for the planning and motor classes: the component
time courses (left), one reconstructed electrode (middle), and the distribution of
component weights (right). 4F–4G where high-weight electrodes sit.

THEIR METHOD (STAR Methods, "Nonnegative matrix factorization"): NMF applied WITHIN the
planning class and WITHIN the motor class separately, not across everything. Factor
count from an RMSE/VAF elbow (Figs S5A–C), roughly three per class — planning splits
early/middle/late, motor splits pre-onset/post-onset/offset.

THE RESULT THAT MATTERS TO US
They clustered electrodes by NMF factor weight — k-means, 1–10 clusters, 1000
iterations — and found no clustering (Figs S5D–S5E). So they never took a partition from
the NMF; weights stay continuous and are plotted as distributions.

Because of that they switched the unit of analysis to the voxel. THAT IS THE PIVOT WE
BORROWED FOR FIG C.11.

TWO THINGS TO TRY FROM THIS SLIDE
  1. Decompose WITHIN a coarse class rather than across all electrodes at once.
  2. Note they apply NO stability filter to their NMF — 100,000 replicates from a good
     seed is an optimiser guarantee, not a data one. Notebook 238 is ahead here."""),

    ("Castellucci S1 — how they decide an electrode responds",
     "Their gate ↔ ours, side by side",
     G3, CA,
     """THEIR METHOD (STAR Methods, "Detection of significant neural responses")
S1A a reaction-time cutoff so trials are behaviourally comparable.
S1B artefact rejection.
S1C the average response against a null built by re-aligning to RANDOM timepoints in the
task period, 1000 times, giving a 95% confidence band.
S1D the same null used for the DURATION of a putative active period.
S1E a period counts only if it is both outside the band and longer than 95% of null
periods.

An electrode is modulated if it shows significant activation OR SUPPRESSION near at
least one of four alignment points. This keeps 1296 of 1784 electrodes — 72.6%.

They note their earlier GLM-based method kept only 32%, and treat the more permissive
detection as the improvement.

OURS — the slide. An amplitude threshold keeping 43%. No null, no duration criterion,
no test of repeatability. One bin in 38,700 separates kept from discarded.

TWO FEATURES OF THEIRS WE DO NOT HAVE
  - a null distribution, so "responsive" means significant rather than large
  - a DURATION criterion, so a brief spike cannot pass

Note their approach needs the continuous signal to re-align at random timepoints, so
like split-half reliability it requires a stage-01 change."""),

    ("What we reproduced",
     "Method by method",
     None, GREEN,
     """FROM NORMAN-HAIGNERE
  Reproduced — the matrix formulation and graded weights (their Fig 1D)
  Reproduced — over-factor then filter, with the model left at K (Figs 1F/1G)
  Reproduced — an initialisation filter and a subject-dominance filter (238)
  Adopting  — split-half reliability as the electrode criterion (Fig 1B)
  Adopting  — noise-corrected variance explained (Fig 1E), once the halves exist
  Cannot do — model-matched stimulus synthesis (Fig 3); needs new stimuli

FROM CASTELLUCCI
  Reproduced — concatenating task epochs into one feature vector per electrode
  Reproduced — comparing several clustering methods over a K range on silhouette
  Reproduced — testing the clustering against a null built from shuffled data
  Reproduced — the voxel permutation test for anatomy (their Fig 3 machinery)
  Reproduced — THE PIVOT: when weights do not cluster, test locations instead
  Not done   — PCA before clustering; they use 3 PCs, we use 900 dimensions
  Not done   — a significance-and-duration response detector (their S1)
  Not done   — decomposing within a coarse class rather than across everything"""),

    ("What we add, and what we still owe",
     "Close here",
     None, CA,
     """WHAT NEITHER PAPER DOES, AND WE DO
  A within-patient shuffle on the component-weight anatomy. Castellucci run this
  control for their class test but not their weight test; Norman-Haignere run no
  anatomical test at all.

  A multiplicity correction on the voxel map. Their voxels step 1 mm apart, denser than
  our 2 mm, reported uncorrected. 746 of our 14,927 voxels pass on noise alone.

  A bootstrap stability filter with a confidence interval and a CALIBRATED threshold.
  Hennig's 0.60 is not a level derived for a given K — on our data the null reaches
  0.775 at K=4, above his "stable" bar.

  Bi-cross-validation holding out rows AND columns.

WHAT WE STILL OWE
  Replace the amplitude gate — machinery built and tested, needs 140 re-run.
  Noise-correct the K curve with the same halves.
  Test whether C.7's negative is dimensionality: cluster in 3 PCs and re-run the null.
  Try decomposing within a coarse class, as Castellucci do.

THE HONEST SUMMARY
We are ahead of both on component stability and on anatomical statistics. We are behind
both on electrode selection. Everything above follows from that one asymmetry."""),
]


def fig(rel):
    p = REPO / rel
    if not p.is_file():
        raise FileNotFoundError(rel)
    return p


def add_slide(prs, idx, title, kicker, image, accent, notes):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    bar = s.shapes.add_shape(1, 0, 0, W, Emu(int(Inches(0.055))))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); bar.shadow.inherit = False

    tb = s.shapes.add_textbox(MARGIN, Inches(0.30), W - 2 * MARGIN, Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = INK
    r.font.name = "Calibri"
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = kicker
    r2.font.size = Pt(13); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"

    if image:
        src = fig(image)
        iw, ih = Image.open(src).size
        top = Inches(1.45)
        avail_w, avail_h = W - 2 * MARGIN, H - top - Inches(0.45)
        sc = min(avail_w / iw, avail_h / ih)
        w, h = int(iw * sc), int(ih * sc)
        s.shapes.add_picture(str(src), int((W - w) / 2),
                             int(top + (avail_h - h) / 2), width=w, height=h)
    else:
        note = s.shapes.add_textbox(MARGIN, Inches(2.6), W - 2 * MARGIN, Inches(1.2))
        nf = note.text_frame; nf.word_wrap = True
        np_ = nf.paragraphs[0]; np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run()
        nr.text = "— speak from the notes —"
        nr.font.size = Pt(15); nr.font.color.rgb = MUTED; nr.font.name = "Calibri"

    n = s.notes_slide.notes_text_frame
    n.text = notes
    for para in n.paragraphs:
        for run in para.runs:
            run.font.size = Pt(11)
            run.font.name = "Consolas" if ("  " in para.text) else "Calibri"

    num = s.shapes.add_textbox(W - Inches(0.95), H - Inches(0.45),
                               Inches(0.6), Inches(0.3))
    tfn = num.text_frame; pn = tfn.paragraphs[0]; pn.alignment = PP_ALIGN.RIGHT
    rn = pn.add_run(); rn.text = str(idx)
    rn.font.size = Pt(11); rn.font.color.rgb = MUTED; rn.font.name = "Consolas"


def main() -> int:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for i, (t, k, img, acc, notes) in enumerate(SLIDES, start=1):
        add_slide(prs, i, t, k, img, acc, notes)
    OUT.mkdir(parents=True, exist_ok=True)
    p = OUT / "FBM_reference_papers.pptx"
    prs.save(str(p))
    print(f"{len(SLIDES)} slides, {sum(1 for s in SLIDES if s[2])} with figures")
    print(f"-> {p}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
