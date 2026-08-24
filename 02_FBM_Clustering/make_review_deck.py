#!/usr/bin/env python3
"""
make_review_deck.py - build the supervisor-meeting deck.

One slide per step of the argument. The bullets go in the SPEAKER NOTES, not on the
slide, so the projected slide is a title and the figure and nothing competes with it
while talking. Every figure is the file the analysis_status page itself points at, so
the deck cannot drift from the site.

Numbers in the notes are the published ones; the sources are named per slide so any of
them can be traced back.

    python make_review_deck.py
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent
OUT = ROOT / "outputs" / "review_deck"

INK = RGBColor(0x1B, 0x23, 0x2C)
MUTED = RGBColor(0x68, 0x72, 0x7D)
ACCENT = RGBColor(0x4A, 0x6F, 0xA5)
RED = RGBColor(0xC1, 0x12, 0x1F)
GREEN = RGBColor(0x1B, 0x78, 0x37)

W, H = Inches(13.333), Inches(7.5)          # 16:9
MARGIN = Inches(0.55)


def fig(rel):
    p = REPO / rel
    if not p.is_file():
        raise FileNotFoundError(rel)
    return p


# (title, kicker, image, accent, notes)
SLIDES = [
    ("Are the response types real?",
     "One month, one question",
     None, ACCENT,
     """THE MONTH IN EIGHT LINES — read this list, then walk it.

1. Published K=7 as a run — convex NMF on 27 patients
2. Tested whether the partition is real — three tests
3. Tested the responsiveness gate — what it removes
4. Replaced the partition with a voxel test
5. Filtered components by stability — over-factor + bootstrap
6. Mapped the language network — LanA / Fedorenko
7. Corrected three published results
8. Rebuilt stage 01 — GO-locked, split-half cubes

Cohort: 27 patients. Gated 1266 electrodes, ungated 2946, localised 1029.
analysis_status is the live appendix — project it rather than rebuilding figures."""),

    ("Definitions",
     "Say once, refer back all meeting",
     None, MUTED,
     """Convex NMF — X ≈ G(W′X), with W, G ≥ 0. Components are convex combinations of the
electrodes themselves, so signed dB is admissible. Ding, Li & Jordan 2010.

Loading w_j — an electrode's share of component j. The seven rows sum to 1.

Silhouette — compactness against the nearest other cluster. NOT space-free; this
matters twice this meeting.

Bootstrap Jaccard — resample electrodes, refit, overlap with the best match.
Hennig 2007: >0.75 stable, 0.60–0.75 a pattern, <0.60 not a cluster.

The gate — n_high_activity > 0. An amplitude threshold, and what the month ends up
questioning."""),

    ("Where we were",
     "FIG C.3 · cnmf/concat_hg",
     "02_FBM_Clustering/outputs/clustering/decomposition/concat_hg/D3_graded_decomposition.png",
     MUTED,
     """Published K=7 as a run on the 27-patient cohort.

Chose K against three criteria — split-half picks 7, bi-cross-validation peaks at 8.

They do NOT agree on concat_rawds: bi-CV 14, LOPO 9, split-half 2. K=7 is defensible
on the track we publish and is not supported by any criterion on the other one.

LAND THIS: K=7 was published. I spent the month trying to break it."""),

    ("Does the model assign electrodes cleanly?",
     "FIG C.5 — every electrode's loadings, grouped by argmax",
     "02_FBM_Clustering/outputs/clustering/cnmf/concat_hg/runs/20260818_112939/decomposition_drivers/D1_membership.png",
     ACCENT,
     """Median top weight 0.43, against 0.14 for an even split.

Only 34% have a majority component.

236 electrodes (19%) have their top two within 0.05 — those labels are coin flips.

Panel D also finds c2 and c4 at r = −0.89, near mirror images.

LAND THIS: graded, not discrete."""),

    ("Does the hard label still recover the features?",
     "FIG C.6 — two routes to the same question",
     "02_FBM_Clustering/outputs/clustering/cnmf/concat_hg/runs/20260818_112939/decomposition_drivers/D2_drivers.png",
     ACCENT,
     """Compares the loading correlation, which never passes through the argmax, against
Cohen's d for the argmax cluster versus every other electrode.

They agree at r = 0.97–0.99 on every component of both feature sets.

LAND THIS: group-level claims survive the hard labelling. Per-electrode membership
claims do not. That is the boundary — state it exactly this way."""),

    ("Is it a separation at all?",
     "FIG C.7 — each method against its own null, in both spaces",
     "02_FBM_Clustering/outputs/clustering/comparison/C7_separation.png",
     ACCENT,
     """Convex NMF: 0.046 (z +1) in dB, 0.094 (z +15) in unit-norm.
K-means: 0.113 (z +10) in dB, −0.004 (z −25) in unit-norm.
Ward: 0.115 (z +15) and −0.010 (z −11).

Weak agreement between methods — ARI 0.21 / 0.19 / 0.49.

WHY: convex NMF unit-norms each electrode before fitting; k-means and Ward use raw dB.
Per-electrode vector norms span 8.96 to 126.13, a 14x range, so in dB the distance is
dominated by amplitude.

LAND THIS: every method separates decisively at home and collapses in the other's
space. They describe different structure — they do not contradict each other.

This figure was published wrong once. Say so here or on the corrections slide."""),

    ("What does the gate remove?",
     "FIG C.9 — same three methods, gate lifted",
     "02_FBM_Clustering/outputs/clustering/comparison/G1_silhouette_by_gate.png",
     ACCENT,
     """Lifting the gate gives 2946 electrodes against 1266.

K-means builds a cluster 74% made of the removed electrodes; Ward 71%.

Convex NMF has none — its most extreme is 69% against a 57% baseline.

Robust at matched K: the most-added cluster stays 74–78% at K=4, 7 and 10, and the
least-added reaches 0.0% at K=7 and K=10 for both hard methods.

LAND THIS: the hard methods in dB are partly separating responsive from
non-responsive — the failure the gate exists to prevent."""),

    ("Are they a different response type?",
     "FIG C.10 — Cohen's d computed twice per cluster",
     "02_FBM_Clustering/outputs/clustering/comparison/G2_drivers_by_gate.png",
     ACCENT,
     """Cohen's d from the gated members only, then from the added members only.

They agree at r = 0.75–0.95 on every cluster of every method.

But the added ones are quieter: mean |HG| 0.392 dB against 0.639, peak 1.58 against
2.51.

LAND THIS: same response type, weaker. Change the criterion, not the cohort.

This is what sends us to split-half reliability — slide 13."""),

    ("Where does each component live?",
     "FIG C.11 — the electrode dropped as the unit of analysis",
     "02_FBM_Clustering/outputs/clustering/component_anatomy/A1_component_anatomy.png",
     ACCENT,
     """Following Castellucci et al. 2026, who hit the identical wall — their electrodes do
not cluster by NMF factor weight either.

Tests the VOXEL instead: for every 1 cm cube holding ≥3 of the 1029 localised
contacts, is the median loading higher than chance allows? Loadings stay continuous —
nothing is partitioned.

14,927 voxels, 10,000 permutations, fsaverage space.

Adds two things neither paper does: a within-patient shuffle, and a multiplicity
correction (746 of 14,927 voxels pass on noise alone).

SURVIVES: c1 — 708 voxels FDR, 100 contacts, peak −56/−30/4, 51% somatomotor.
          c5 — 242 voxels FDR, 111 contacts, peak −70/−34/−18, 81% default.
FAILS:    c3 and c6 localise nowhere. c2 (n=12) and c4 (n=3) rest on too few contacts.

LAND THIS: two components out of seven have a cortical home."""),

    ("Which components are stable?",
     "Notebook 238 — over-factor, then filter",
     "02_FBM_Clustering/outputs/clustering/overfactor_stability/boot_cnmf_concat_hg_K10.png",
     ACCENT,
     """Norman-Haignere's strategy, with the bootstrap as the primary filter instead of
initialisation. The MODEL stays at K — only interpretation is restricted.

Hennig's 0.60 is a rule of thumb, not a level derived for a given K, n and set of
cluster sizes, so it is calibrated against structureless data with the same covariance:
0.775 at K=4, 0.502 at K=7, 0.344 at K=10.

Margin over chance therefore RISES with K: +0.055, +0.253, +0.341.

At K=7: c5 0.86, c1 0.84, c6 0.82 stable; c4 0.78 and c3 0.76 a pattern; c0 unresolved,
c2 at chance.

LAND THIS: at K=4 the null already reaches 0.775 — above Hennig's "stable" bar. K=4 is
the weakest of the three, not the strongest."""),

    ("LanA — mapping onto the language network",
     "FIG 3a.R1 · six atlas partitions built as clustering runs",
     "02_FBM_Clustering/outputs/clustering/atlas/lana_all/lana_cards.png",
     ACCENT,
     """Built the LanA in/out split as six clustering runs so the project glassbrain renders
it — 2 cohorts x (P ≥ 0.05, P ≥ 0.10, 5-band continuous).

Read DOWN the column and the threshold's effect is visible without the table:
  P ≥ 0.05 — almost everything is "in network", so the cut separates nothing
  P ≥ 0.10 — becomes selective
  top band (P ≥ 0.35) — collapses onto perisylvian temporal cortex, the atlas core

R1 is all electrodes (n=2981); R2 is the concatenated cohort (n=2617), directly
comparable row for row.

Also: multi-label attribution — 70% of Neurosynth-loaded contacts cleanly
attributable, 30% multi-load."""),

    ("ERSP x P(language) — the correlation map",
     "FIG 9 · make_fedorenko_corr.py",
     "04_FBM_Pooling/outputs/pooling/atlas_corr/fedorenko/fig9_corr_map_berlin.png",
     ACCENT,
     """2724 contacts. 26.3% of bins survive FDR.

POSITIVE during the spoken prompt (audio stimulus, ~50–150 Hz) — closer to the
language network means more high gamma while hearing speech.

POSITIVE in all three response windows (~50–250 Hz) — language sites engage during
production regardless of how the trial was cued.

NEGATIVE during visual stimulus encoding (picture & reading, ~200–400 Hz).

LAND THIS: the language network tracks linguistic processing — hearing speech and
producing it — and actively does not track visual sensory encoding. That is a
dissociation, not a coverage effect."""),

    ("Three published results, corrected",
     "All three dated and withdrawn in place on analysis_status",
     None, RED,
     """1. FIG C.7 scored all three methods in raw dB. Convex NMF unit-norms each electrode
first, so it was measured in a space it never optimised in. Corrected: z goes from +1
to +15 in unit-norm.

2. FIG C.9 called the silhouette "the tell". That is a K=4 artefact — a wash at K=7 and
REVERSED at K=10. Withdrawn; the composition claim replaced it and does hold at
K=4, 7 and 10.

3. The bootstrap could not exceed 1 − 1/e = 0.632. Out-of-bag rows sat in the Jaccard
denominator and never in the numerator. Verified rather than argued: feeding it a
perfect recovery returned 0.6322 ± 0.0084 instead of 1.0.
   Result reversed — 0 of 7 clusters surviving became 3 stable, 2 patterns.

LAND THIS: each was caught by a check built to catch it — a matched-K comparison, a
null model, a known-answer test. The tell on the third was "best J = 0.59" at every
single K; a suspiciously flat ceiling is what made me look."""),

    ("Stage 01 rebuilt",
     "Notebooks 140 · 150 · lf_ersp",
     None, MUTED,
     """Added GO-locked real time (notebook 150) — nothing time-warped, epochs cut in real
seconds and centred on the GO cue.

Validated on synthetic data: a burst planted 0.25 s after GO is recovered within one
grid step regardless of window or stimulus length, while onset alignment smears it to
+1.56 s.

Added odd/even half-cubes for split-half reliability — avg_db is bit-identical to the
pre-change code, a real repeatable response scores r = +0.908, a pure-noise channel
r = +0.037.

Found silent failures: duplicate trial TSVs double-weighting every MicroEPI average;
auxiliary channels ERSP'd as cortex; a dash-blind non-neural filter.

Recovered EL046 and EL048 — their lookup directories pointed one level above the
workbooks."""),

    ("What the gate is actually deciding",
     "Borderline contacts, both sides of n_high_activity",
     "02_FBM_Clustering/outputs/clustering/gate_examples/G3_gate_borderline_examples.png",
     RED,
     """THE RULE. For each electrode-condition ERSP (129 freq x 300 time, 0–400 Hz, dB):
  prop_above_pos = fraction of bins > +2.2 dB
  prop_below_neg = fraction of bins < −3.0 dB
  high_activity  = (prop_above_pos ≥ 0.02) OR (prop_below_neg ≥ 0.04)
n_high_activity counts the conditions that pass; keep if ≥1. 1323 of 3002 pass.

Note it uses the WHOLE 0–400 Hz cube — it is not a high-gamma measure.

Each cube holds 129 x 300 = 38,700 bins, so the 2% rule is a count: 774 bins.

EL033 aH_R5 has 774 bins and is KEPT.
EL033 aH_R2 — the SAME SHANK IN THE SAME PATIENT — has 773 and is DISCARDED.

One bin in 38,700 decides it, and the two are indistinguishable by eye.

PAT_3455 HAD5 is kept on the suppression criterion instead (1548 bins below −3 dB),
which shows the OR is doing real work — a suppressed contact can pass without any
positive response at all.

LAND THIS: nothing in this rule measures whether the response REPEATS. That is the
argument for replacing the criterion with split-half reliability rather than moving
the threshold."""),

    ("The claim, and what's next",
     "Close here",
     None, GREEN,
     """NEXT
Replace the amplitude gate with split-half reliability — both reference papers select
on reproducibility. Machinery built and tested; needs 140 re-run.

Noise-correct bi-cross-validation with the same halves — turns K from a relative
comparison into a fraction of EXPLAINABLE variance (Norman-Haignere Fig 1E).

Test whether C.7's negative is dimensionality — Castellucci cluster in 3 PCs, we
cluster in 900.

Clear 736 stale ERSP cubes; recover PAT_3965 and EL043, whose trigger TSVs are missing
for some conditions.

CLOSING LINE
Not seven response types. A K=7 decomposition that generalises to held-out data, three
of whose components are bootstrap-stable, two of which have a reproducible cortical
home — and a cohort criterion I now know is wrong and have the machinery to replace."""),
]


def add_slide(prs, idx, title, kicker, image, accent, notes):
    s = prs.slides.add_slide(prs.slide_layouts[6])          # blank

    bar = s.shapes.add_shape(1, 0, 0, W, Emu(int(Inches(0.055))))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background(); bar.shadow.inherit = False

    tb = s.shapes.add_textbox(MARGIN, Inches(0.30), W - 2 * MARGIN, Inches(0.95))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = INK
    r.font.name = "Calibri"
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = kicker
    r2.font.size = Pt(13); r2.font.color.rgb = MUTED; r2.font.name = "Calibri"

    if image:
        src = fig(image)
        iw, ih = Image.open(src).size
        top = Inches(1.45)
        avail_w = W - 2 * MARGIN
        avail_h = H - top - Inches(0.45)
        scale = min(avail_w / iw, avail_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        s.shapes.add_picture(str(src), int((W - w) / 2), int(top + (avail_h - h) / 2),
                             width=w, height=h)

    n = s.notes_slide.notes_text_frame
    n.text = notes
    for para in n.paragraphs:
        for run in para.runs:
            run.font.size = Pt(12)

    num = s.shapes.add_textbox(W - Inches(0.95), H - Inches(0.45),
                               Inches(0.6), Inches(0.3))
    tfn = num.text_frame; pn = tfn.paragraphs[0]; pn.alignment = PP_ALIGN.RIGHT
    rn = pn.add_run(); rn.text = str(idx)
    rn.font.size = Pt(11); rn.font.color.rgb = MUTED; rn.font.name = "Consolas"
    return s


def main() -> int:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for i, (t, k, img, acc, notes) in enumerate(SLIDES, start=1):
        add_slide(prs, i, t, k, img, acc, notes)
    OUT.mkdir(parents=True, exist_ok=True)
    p = OUT / "FBM_month_review.pptx"
    prs.save(str(p))
    n_img = sum(1 for s in SLIDES if s[2])
    print(f"{len(SLIDES)} slides, {n_img} with figures")
    print(f"-> {p}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
